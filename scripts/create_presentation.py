#!/usr/bin/env python3
"""Build the eight-slide Doggo project presentation as an editable PPTX."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "deliverables" / "Doggo_ROS2_YOLOv10_绝缘子巡检_8页汇报.pptx"
REAL_EVIDENCE = (
    ROOT / "inspection_results" / "defect_20260721_182434_006.jpg"
)
SIM_IMAGE = ROOT / "datasets" / "insulator_sim" / "images" / "val" / "sim_000000.jpg"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Microsoft YaHei"
MONO = "Consolas"

NAVY = "071B33"
NAVY_2 = "0D2B4D"
BLUE = "0B78E3"
CYAN = "11BFD0"
RED = "E9574F"
AMBER = "F3AA35"
GREEN = "28A87A"
INK = "17324D"
MUTED = "607489"
PALE = "EAF2F8"
PALE_BLUE = "E9F4FF"
PALE_CYAN = "E9FAFB"
PALE_RED = "FFF0EF"
PALE_AMBER = "FFF6E8"
WHITE = "FFFFFF"
LIGHT = "F6F9FC"
LINE = "D7E3ED"


def color(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 16,
    font_color: str = INK,
    bold: bool = False,
    font: str = FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
    line_spacing: float = 1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    parts = text.split("\n")
    for index, part in enumerate(parts):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = part
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color(font_color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = WHITE,
    line: str | None = None,
    radius: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    if line:
        shape.line.color.rgb = color(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except (IndexError, ValueError):
            pass
    return shape


def add_line(slide, x1, y1, x2, y2, line_color=LINE, width=1.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color(line_color)
    line.line.width = Pt(width)
    return line


def add_picture_fill(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        image_w, image_h = image.size
    box_ratio = w / h
    image_ratio = image_w / image_h
    picture = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        visible = image_ratio / box_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_image_frame(slide, path: Path, x, y, w, h, caption: str):
    add_rect(slide, x - 0.03, y - 0.03, w + 0.06, h + 0.38, WHITE, LINE, True)
    add_picture_fill(slide, path, x, y, w, h)
    add_text(slide, caption, x + 0.08, y + h + 0.07, w - 0.16, 0.22, 9.5, MUTED)


def add_power_motif(slide, dark=False):
    base = "315B7C" if dark else "B8D2E5"
    accent = CYAN if dark else BLUE
    points = [(10.65, 0.32), (11.35, 0.20), (12.10, 0.42), (12.78, 0.22)]
    for first, second in zip(points, points[1:]):
        add_line(slide, first[0], first[1], second[0], second[1], base, 1.2)
    for index, (x, y) in enumerate(points):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.055), Inches(y - 0.055), Inches(0.11), Inches(0.11)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color(accent if index in (0, 3) else base)
        circle.line.fill.background()


def add_header(slide, section: str, title: str, number: int):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT)
    add_rect(slide, 0, 0, 0.12, SLIDE_H, BLUE)
    add_text(slide, section.upper(), 0.55, 0.30, 3.4, 0.27, 10, BLUE, True, MONO)
    add_text(slide, title, 0.52, 0.62, 10.7, 0.58, 25, NAVY, True)
    add_text(slide, f"0{number}", 12.15, 0.62, 0.62, 0.34, 11, MUTED, True, MONO, PP_ALIGN.RIGHT)
    add_line(slide, 0.55, 1.25, 12.78, 1.25, LINE, 1.0)
    add_power_motif(slide)


def add_footer(slide, number: int, note: str = "Doggo · ROS 2 + YOLOv10 输电巡检"):
    add_line(slide, 0.55, 7.08, 12.78, 7.08, LINE, 0.8)
    add_text(slide, note, 0.55, 7.14, 7.0, 0.18, 8.5, MUTED, False, FONT)
    add_text(slide, f"{number} / 8", 11.85, 7.14, 0.9, 0.18, 8.5, MUTED, True, MONO, PP_ALIGN.RIGHT)


def add_badge(slide, text, x, y, w, fill=PALE_BLUE, font_color=BLUE):
    add_rect(slide, x, y, w, 0.34, fill, None, True)
    add_text(slide, text, x + 0.08, y + 0.03, w - 0.16, 0.24, 9.5, font_color, True, FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def add_card(slide, x, y, w, h, title, body, accent=BLUE, number=None):
    add_rect(slide, x, y, w, h, WHITE, LINE, True)
    add_rect(slide, x, y, 0.07, h, accent, None, True)
    if number:
        add_text(slide, number, x + 0.23, y + 0.18, 0.38, 0.28, 10, accent, True, MONO)
        title_x = x + 0.72
    else:
        title_x = x + 0.27
    add_text(slide, title, title_x, y + 0.13, w - (title_x - x) - 0.18, 0.35, 14, NAVY, True)
    add_text(slide, body, x + 0.27, y + 0.62, w - 0.49, h - 0.78, 10.5, MUTED, False, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 0.02, 1.0)


def add_metric(slide, x, y, w, value, label, accent=BLUE):
    add_rect(slide, x, y, w, 0.92, WHITE, LINE, True)
    add_text(slide, value, x + 0.12, y + 0.09, w - 0.24, 0.42, 23, accent, True, MONO)
    add_text(slide, label, x + 0.12, y + 0.57, w - 0.24, 0.20, 9.5, MUTED, False)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_picture_fill(slide, REAL_EVIDENCE, 7.42, 0, 5.913, 7.5)
    add_rect(slide, 7.18, 0, 0.34, 7.5, CYAN)
    add_rect(slide, 7.03, 0, 0.15, 7.5, BLUE)
    add_text(slide, "POWER · AI · ENGINEERING", 0.66, 0.48, 4.2, 0.26, 10, CYAN, True, MONO)
    add_text(slide, "面向新型电力系统的", 0.66, 1.22, 5.9, 0.50, 25, WHITE, False)
    add_text(slide, "绝缘子缺陷智能巡检系统", 0.66, 1.73, 6.1, 0.72, 31, WHITE, True)
    add_text(slide, "ROS 2 + YOLOv10s · 仿真到真实场景", 0.70, 2.67, 5.7, 0.38, 15, "A9C7DE")
    add_line(slide, 0.70, 3.33, 5.92, 3.33, "315B7C", 1.2)
    for i, (label, x, width) in enumerate(
        [("能源电力", 0.70, 1.28), ("智能感知", 2.13, 1.28), ("工程实践", 3.56, 1.28)]
    ):
        add_badge(slide, label, x, 3.62, width, NAVY_2, CYAN if i == 1 else WHITE)
    add_text(
        slide,
        "8 页项目成果汇报\n可训练 · 可运行 · 可追溯 · 能说明边界",
        0.70,
        4.55,
        5.8,
        0.95,
        13,
        "C8D9E7",
        False,
        FONT,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.TOP,
        0.02,
        1.2,
    )
    add_text(slide, "2026 · Doggo Project", 0.70, 6.75, 3.3, 0.25, 9.5, "7997B0", True, MONO)
    add_rect(slide, 8.02, 6.65, 4.70, 0.48, NAVY)
    add_text(
        slide,
        "真实推理证据｜missing_disc 0.91",
        8.18,
        6.75,
        4.35,
        0.20,
        9.5,
        WHITE,
        True,
        FONT,
        PP_ALIGN.RIGHT,
    )


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "POWER-SYSTEM CONTEXT", "从线路安全痛点到工程目标", 2)
    add_text(
        slide,
        "电力高校特色：让电气设备机理、人工智能与机器人系统在同一个工程闭环中相遇。",
        0.58,
        1.42,
        11.8,
        0.40,
        13,
        MUTED,
    )
    cards = [
        ("高风险、长距离", "高压线路人工复核成本高，山地与跨区通道覆盖困难。", RED),
        ("小目标、强背景", "缺片区域远小于绝缘子串，杆塔、导线和天空容易形成干扰。", AMBER),
        ("故障少、数据贵", "真实缺陷样本稀缺，跨线路、相机与天气的域偏移明显。", BLUE),
    ]
    for index, (title, body, accent) in enumerate(cards):
        add_card(slide, 0.58, 2.02 + index * 1.25, 5.05, 1.05, title, body, accent, f"0{index + 1}")
    add_rect(slide, 6.05, 2.02, 6.70, 3.55, NAVY, None, True)
    add_text(slide, "工程目标", 6.43, 2.30, 2.0, 0.35, 12, CYAN, True, MONO)
    goals = [
        ("看得见", "定位真实输电背景中的绝缘子串"),
        ("看得准", "识别 missing_disc 缺片并给出置信度"),
        ("交得出", "ROS 2 告警、证据与可复现权重同步交付"),
    ]
    for index, (title, body) in enumerate(goals):
        y = 2.92 + index * 0.72
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.43), Inches(y), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color(CYAN if index == 1 else BLUE)
        circle.line.fill.background()
        add_text(slide, str(index + 1), 6.43, y + 0.06, 0.42, 0.20, 9.5, WHITE, True, MONO, PP_ALIGN.CENTER)
        add_text(slide, title, 7.05, y - 0.01, 1.15, 0.30, 13, WHITE, True)
        add_text(slide, body, 8.22, y - 0.01, 4.15, 0.35, 10.5, "BFD3E3")
    add_rect(slide, 0.58, 5.90, 12.17, 0.77, PALE_BLUE, None, True)
    add_text(slide, "项目定位", 0.85, 6.10, 1.10, 0.24, 11, BLUE, True)
    add_text(slide, "面向课程实践与算法原型的可运行巡检系统，而非未经现场验收的安全决策系统。", 1.95, 6.08, 10.35, 0.30, 12.5, NAVY, True)
    add_footer(slide, 2)


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "FULL TOOLCHAIN", "完整工具栈与数据闭环", 3)
    add_text(slide, "训练链路与运行链路解耦，但共享同一 ROS 2 感知接口。", 0.58, 1.42, 9.5, 0.35, 13, MUTED)
    stages = [
        ("01", "Gazebo", "自动采图\n3D→2D 标注", BLUE),
        ("02", "YOLOv10s", "仿真预训练\n500 张 / 30 epoch", CYAN),
        ("03", "CPLID", "真实背景迁移\n678 train / 170 val", AMBER),
        ("04", "ROS 2", "相机 / 图片 / 视频\n统一推理节点", BLUE),
        ("05", "巡检闭环", "告警 + JPG + JSON\n结果可复核", GREEN),
    ]
    start_x = 0.58
    card_w = 2.18
    gap = 0.24
    for index, (number, title, body, accent) in enumerate(stages):
        x = start_x + index * (card_w + gap)
        add_rect(slide, x, 2.12, card_w, 2.35, WHITE, LINE, True)
        add_rect(slide, x, 2.12, card_w, 0.13, accent, None, True)
        add_text(slide, number, x + 0.18, 2.42, 0.45, 0.24, 10, accent, True, MONO)
        add_text(slide, title, x + 0.18, 2.77, card_w - 0.36, 0.38, 15, NAVY, True)
        add_text(slide, body, x + 0.18, 3.38, card_w - 0.36, 0.68, 10.3, MUTED, False, FONT, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 0.01, 1.12)
        if index < len(stages) - 1:
            arrow_x = x + card_w + 0.045
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(arrow_x), Inches(3.05), Inches(0.18), Inches(0.45))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color("9FC3DD")
            arrow.line.fill.background()
    add_text(slide, "与业界前沿工程范式一致", 0.60, 4.92, 3.1, 0.30, 12, NAVY, True)
    trends = [
        ("SYNTHETIC DATA", "仿真扩充稀缺样本"),
        ("SIM-TO-REAL", "跨域迁移而非直接套用"),
        ("EDGE-AWARE", "轻量模型与低延迟推理"),
        ("MLOPS TRACE", "权重、指标、校验值留痕"),
    ]
    for index, (tag, body) in enumerate(trends):
        x = 0.60 + index * 3.05
        add_rect(slide, x, 5.42, 2.78, 0.95, PALE_CYAN if index % 2 else PALE_BLUE, None, True)
        add_text(slide, tag, x + 0.15, 5.56, 2.48, 0.20, 9, BLUE if index % 2 == 0 else GREEN, True, MONO)
        add_text(slide, body, x + 0.15, 5.88, 2.48, 0.25, 10.5, NAVY, True)
    add_footer(slide, 3)


def slide_simulation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "DIGITAL POWER LAB", "仿真数据引擎：把电力设备先带进数字世界", 4)
    add_image_frame(slide, SIM_IMAGE, 0.58, 1.60, 5.55, 4.37, "实际 Gazebo 训练帧｜绝缘子姿态、距离与视角可自动变化")
    add_rect(slide, 6.48, 1.60, 6.28, 1.42, NAVY, None, True)
    add_text(slide, "自动标注原理", 6.78, 1.86, 2.1, 0.28, 12, CYAN, True, MONO)
    add_text(slide, "3D 包围盒  →  相机内外参  →  2D YOLO 框", 6.78, 2.24, 5.40, 0.36, 17, WHITE, True, MONO)
    add_text(slide, "避免逐帧人工标注，并保证仿真标签与模型位姿一致。", 6.78, 2.65, 5.35, 0.24, 10.5, "BCD1E1")
    factors = [
        ("几何域", "距离、平移、俯仰/横滚/偏航"),
        ("成像域", "亮度、色彩、模糊、噪声"),
        ("样本域", "正样本 + 少量负样本"),
    ]
    for index, (title, body) in enumerate(factors):
        y = 3.30 + index * 0.76
        add_rect(slide, 6.48, y, 6.28, 0.59, WHITE, LINE, True)
        add_text(slide, title, 6.75, y + 0.13, 1.02, 0.25, 11, BLUE, True)
        add_text(slide, body, 7.92, y + 0.13, 4.44, 0.25, 10.5, NAVY, True)
    add_metric(slide, 6.48, 5.73, 1.84, "500", "仿真训练图片", BLUE)
    add_metric(slide, 8.48, 5.73, 1.84, "30", "训练 epoch", CYAN)
    add_metric(slide, 10.48, 5.73, 2.28, "0.976", "仿真 mAP50", GREEN)
    add_footer(slide, 4)


def slide_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "MODEL DESIGN", "模型设计：面向小缺片的双域优化", 5)
    add_text(slide, "不追求一个权重包打天下，而是保留不同视觉域的最佳模型。", 0.58, 1.42, 10.0, 0.35, 13, MUTED)
    add_rect(slide, 0.58, 2.03, 2.62, 1.52, PALE_BLUE, None, True)
    add_text(slide, "YOLOv10s", 0.86, 2.30, 2.06, 0.40, 20, BLUE, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, "轻量、面向端到端与低延迟", 0.85, 2.88, 2.10, 0.30, 10.5, MUTED, False, FONT, PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(3.48), Inches(2.49), Inches(0.60), Inches(0.58))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color(CYAN); arrow.line.fill.background()
    add_rect(slide, 4.34, 1.75, 3.62, 1.52, WHITE, BLUE, True)
    add_text(slide, "仿真域权重", 4.64, 2.02, 1.20, 0.30, 12, BLUE, True)
    add_text(slide, "insulator_yolov10s.pt", 4.64, 2.48, 2.96, 0.28, 12, NAVY, True, MONO)
    add_text(slide, "1 类 · 640 px", 4.64, 2.86, 2.90, 0.24, 10, MUTED)
    add_rect(slide, 4.34, 3.55, 3.62, 1.52, WHITE, GREEN, True)
    add_text(slide, "真实域缺陷权重", 4.64, 3.82, 1.70, 0.30, 12, GREEN, True)
    add_text(slide, "insulator_defect_yolov10s.pt", 4.64, 4.28, 3.00, 0.28, 11, NAVY, True, MONO)
    add_text(slide, "2 类 · 768 px · 缺片小目标", 4.64, 4.66, 2.95, 0.24, 10, MUTED)
    add_line(slide, 6.15, 3.27, 6.15, 3.55, CYAN, 2.4)
    add_text(slide, "SIM-TO-REAL", 6.35, 3.29, 1.30, 0.20, 8.5, CYAN, True, MONO)
    decisions = [
        ("768 px", "保留远距离航拍中的小缺片纹理", BLUE),
        ("双类别", "绝缘子整体框 + 缺片局部框", GREEN),
        ("双权重", "降低真实域微调造成的仿真遗忘", AMBER),
        ("可追溯", "训练配置、指标与 SHA-256 写入模型卡", RED),
    ]
    for index, (title, body, accent) in enumerate(decisions):
        y = 1.75 + index * 1.02
        add_rect(slide, 8.35, y, 4.40, 0.81, WHITE, LINE, True)
        add_rect(slide, 8.35, y, 0.09, 0.81, accent, None, True)
        add_text(slide, title, 8.68, y + 0.15, 1.05, 0.25, 11.5, accent, True, MONO)
        add_text(slide, body, 9.75, y + 0.15, 2.70, 0.35, 10.3, NAVY, True)
    add_rect(slide, 0.58, 5.62, 12.17, 0.86, NAVY, None, True)
    add_text(slide, "前沿一致性", 0.86, 5.87, 1.25, 0.24, 11, CYAN, True)
    add_text(slide, "端到端检测 × 合成数据 × 迁移学习 × ROS 2 模块化 × 事件级证据", 2.20, 5.82, 9.90, 0.33, 14, WHITE, True)
    add_footer(slide, 5)


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "MODEL ARTIFACTS", "真实模型产物与可量化结果", 6)
    add_image_frame(slide, REAL_EVIDENCE, 0.58, 1.58, 7.05, 4.73, "实际 ROS 2 推理证据｜真实输电背景；CPLID 缺陷目标为合成缺片")
    add_metric(slide, 7.98, 1.58, 2.18, "0.980", "missing_disc mAP50", GREEN)
    add_metric(slide, 10.37, 1.58, 2.38, "10.7 ms", "RTX 2080 Ti 单帧", BLUE)
    add_metric(slide, 7.98, 2.67, 2.18, "0.914", "缺片 Precision", CYAN)
    add_metric(slide, 10.37, 2.67, 2.38, "0.960", "缺片 Recall", AMBER)
    add_rect(slide, 7.98, 3.85, 4.77, 1.19, NAVY, None, True)
    add_text(slide, "可交付权重", 8.24, 4.05, 1.22, 0.22, 10, CYAN, True, MONO)
    add_text(slide, "insulator_defect_yolov10s.pt", 8.24, 4.39, 4.12, 0.25, 11.3, WHITE, True, MONO)
    add_text(slide, "16.53 MB · SHA-256  a61f1bea…c273f6", 8.24, 4.72, 4.12, 0.22, 9, "B8CCDC", False, MONO)
    add_rect(slide, 7.98, 5.25, 4.77, 1.43, WHITE, LINE, True)
    add_text(slide, "两类平均", 8.24, 5.45, 1.15, 0.22, 10, BLUE, True)
    add_text(slide, "P 0.791   R 0.844", 9.55, 5.45, 2.72, 0.22, 10.5, NAVY, True, MONO)
    add_text(slide, "mAP50", 8.24, 5.83, 1.15, 0.22, 10, BLUE, True)
    add_text(slide, "0.882", 9.55, 5.80, 1.20, 0.27, 14, NAVY, True, MONO)
    add_text(slide, "mAP50-95", 10.55, 5.83, 1.15, 0.22, 10, BLUE, True)
    add_text(slide, "0.530", 11.58, 5.80, 0.82, 0.27, 14, NAVY, True, MONO, PP_ALIGN.RIGHT)
    add_footer(slide, 6)


def slide_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "ROS 2 DELIVERY", "工程交付：检测不是终点", 7)
    sources = ["Gazebo 相机", "真实 ROS 相机", "图片 / 目录", "视频文件"]
    outputs = ["/doggo/detections", "/doggo/defect_alerts", "/doggo/detection_image", "JPG + JSON 证据"]
    add_text(slide, "输入", 0.58, 1.54, 1.0, 0.25, 11, BLUE, True, MONO)
    for index, source in enumerate(sources):
        y = 1.92 + index * 0.72
        add_rect(slide, 0.58, y, 2.57, 0.53, WHITE, LINE, True)
        add_text(slide, source, 0.80, y + 0.13, 2.12, 0.22, 10.5, NAVY, True)
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(3.40), Inches(2.67), Inches(0.55), Inches(0.76))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color(CYAN); arrow.line.fill.background()
    add_rect(slide, 4.17, 1.85, 4.22, 2.64, NAVY, None, True)
    add_text(slide, "YOLOv10 ROS 2 NODE", 4.52, 2.23, 3.52, 0.30, 12, CYAN, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, "统一推理节点", 4.52, 2.80, 3.52, 0.38, 20, WHITE, True, FONT, PP_ALIGN.CENTER)
    add_text(slide, "模型切换 · 阈值配置 · 证据冷却", 4.52, 3.44, 3.52, 0.28, 10.5, "BFD2E2", False, FONT, PP_ALIGN.CENTER)
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(8.64), Inches(2.67), Inches(0.55), Inches(0.76))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = color(CYAN); arrow.line.fill.background()
    add_text(slide, "输出", 9.45, 1.54, 1.0, 0.25, 11, GREEN, True, MONO)
    for index, output in enumerate(outputs):
        y = 1.92 + index * 0.72
        add_rect(slide, 9.45, y, 3.30, 0.53, WHITE, LINE, True)
        add_text(slide, output, 9.67, y + 0.13, 2.86, 0.22, 10.2, NAVY, True, MONO if output.startswith("/") else FONT)
    add_rect(slide, 0.58, 5.10, 7.35, 1.26, NAVY, None, True)
    add_text(slide, "一键运行", 0.86, 5.32, 1.0, 0.22, 10, CYAN, True, MONO)
    add_text(slide, "./scripts/run_sim.sh", 0.86, 5.72, 2.78, 0.24, 11, WHITE, True, MONO)
    add_text(slide, "ros2 launch ... real_image_demo.launch.py source:=...", 3.78, 5.72, 3.80, 0.24, 9.2, WHITE, True, MONO)
    add_rect(slide, 8.18, 5.10, 4.57, 1.26, PALE_BLUE, None, True)
    add_text(slide, "GitHub 规范产物", 8.45, 5.32, 1.95, 0.22, 10, BLUE, True, MONO)
    add_text(slide, "README · 模型卡 · MIT · CI", 8.45, 5.72, 3.80, 0.24, 11, NAVY, True)
    add_text(slide, "Git LFS · Issue / PR 模板", 8.45, 6.02, 3.80, 0.22, 9.5, MUTED)
    add_footer(slide, 7)


def slide_limits(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "LIMITATIONS & ROADMAP", "现阶段不足：可交付，也要把边界说清楚", 8)
    add_text(slide, "工程严谨性不只是给出高指标，更是说明指标在什么条件下成立。", 0.58, 1.42, 10.8, 0.35, 13, MUTED)
    headers = [("现阶段不足", 0.58, 3.06), ("可能影响", 3.72, 3.03), ("下一阶段改进", 6.85, 5.90)]
    for title, x, w in headers:
        add_rect(slide, x, 1.92, w, 0.49, NAVY, None, True)
        add_text(slide, title, x + 0.16, 2.04, w - 0.32, 0.22, 10.5, WHITE, True)
    rows = [
        ("仅有 missing_disc 一类缺陷", "无法覆盖裂纹、污损、闪络等故障", "建立多缺陷标签体系；引入分割/异常检测"),
        ("CPLID 缺片目标为合成样本", "真实故障纹理与背景融合仍有域差异", "采集真实故障；电力专家复核标注"),
        ("指标来自同源固定划分", "可能高估跨地区、相机和天气泛化", "开展跨线路、跨设备外部测试"),
        ("尚未完成现场飞行验收", "阈值、漏检与误报成本未被充分量化", "建立故障案例库与分级验收流程"),
        ("仅验证桌面 GPU 推理", "Jetson/边缘终端的功耗与实时性未知", "ONNX/TensorRT、量化与端侧压力测试"),
    ]
    row_y = 2.51
    for index, row in enumerate(rows):
        fill = WHITE if index % 2 == 0 else "F0F5F9"
        add_rect(slide, 0.58, row_y, 12.17, 0.65, fill, LINE, False)
        add_text(slide, row[0], 0.76, row_y + 0.12, 2.70, 0.36, 10.2, NAVY, True)
        add_text(slide, row[1], 3.90, row_y + 0.12, 2.70, 0.38, 9.8, MUTED)
        add_text(slide, row[2], 7.05, row_y + 0.12, 5.38, 0.38, 9.8, BLUE, True)
        row_y += 0.67
    add_rect(slide, 0.58, 6.12, 12.17, 0.59, PALE_RED, None, True)
    add_text(slide, "当前定位", 0.82, 6.28, 1.02, 0.22, 10.5, RED, True)
    add_text(slide, "课程成果 / 算法原型 / 现场数据预筛工具；不能替代电力安全专业复核。", 1.93, 6.25, 10.35, 0.27, 12, NAVY, True)
    add_footer(slide, 8, "Doggo · 结论：工程可运行，指标可追溯，边界可解释")


def set_core_properties(prs):
    props = prs.core_properties
    props.title = "Doggo ROS 2 + YOLOv10 绝缘子缺陷智能巡检系统"
    props.subject = "电力高校特色的八页项目成果汇报"
    props.author = "Doggo Project"
    props.keywords = "ROS 2, YOLOv10, 绝缘子, 输电线路, 缺陷检测, Sim-to-Real"
    props.comments = "Metrics and limitations are documented from the local project artifacts."


def main():
    for required in (REAL_EVIDENCE, SIM_IMAGE):
        if not required.is_file():
            raise SystemExit(f"Required presentation asset not found: {required}")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    set_core_properties(prs)
    slide_cover(prs)
    slide_problem(prs)
    slide_pipeline(prs)
    slide_simulation(prs)
    slide_model(prs)
    slide_results(prs)
    slide_engineering(prs)
    slide_limits(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Presentation generated: {OUTPUT}")


if __name__ == "__main__":
    main()
