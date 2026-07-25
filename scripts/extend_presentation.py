"""Extend the InsuLens presentation with YOLO fine-tuning material.

The script keeps the original deck untouched and writes a 12-slide derivative.
It deliberately reuses the original slide chrome so the added pages remain
visually consistent with the existing report.
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "reports" / "InsuLens.pptx"
OUTPUT = ROOT / "reports" / "InsuLens_12pages.pptx"

FONT = "宋体"
MONO = "Consolas"

NAVY = "071D33"
NAVY_2 = "0B2947"
TEXT = "0A2035"
MUTED = "63798C"
LIGHT_MUTED = "AFC0CE"
LINE = "D3DFE8"
BLUE = "147DDF"
CYAN = "19C6D4"
GREEN = "23AE86"
ORANGE = "F3A11A"
RED = "EF5B5B"
PURPLE = "7A78E8"
WHITE = "FFFFFF"
PALE_BLUE = "EAF4FC"
PALE_CYAN = "E8F8F8"
PALE_GREEN = "EAF7F2"
PALE_ORANGE = "FFF5E5"
PALE_RED = "FDEEEE"
PALE_PURPLE = "F0EFFF"
PALE_GRAY = "F1F6FA"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str | None) -> None:
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str | None = LINE, width: float = 1.0) -> None:
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)


def format_text_frame(
    text_frame,
    *,
    size: float,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    font: str = FONT,
    margins: tuple[float, float, float, float] = (0.05, 0.05, 0.03, 0.03),
) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = valign
    text_frame.margin_left = Inches(margins[0])
    text_frame.margin_right = Inches(margins[1])
    text_frame.margin_top = Inches(margins[2])
    text_frame.margin_bottom = Inches(margins[3])
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: float = 10.0,
    color: str = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    font: str = FONT,
    margins: tuple[float, float, float, float] = (0.02, 0.02, 0.02, 0.02),
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    format_text_frame(
        shape.text_frame,
        size=size,
        color=color,
        bold=bold,
        align=align,
        valign=valign,
        font=font,
        margins=margins,
    )
    shape.text_frame.paragraphs[0].runs[0].text = text
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str = "",
    *,
    fill: str = WHITE,
    line: str | None = LINE,
    radius: bool = True,
    size: float = 10.0,
    color: str = TEXT,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    margins: tuple[float, float, float, float] = (0.08, 0.08, 0.05, 0.05),
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line)
    if text:
        format_text_frame(
            shape.text_frame,
            size=size,
            color=color,
            bold=bold,
            align=align,
            valign=valign,
            margins=margins,
        )
        shape.text_frame.paragraphs[0].runs[0].text = text
    return shape


def add_badge(slide, x: float, y: float, text: str, color: str, size: float = 9.0):
    return add_box(
        slide,
        x,
        y,
        0.36,
        0.36,
        text,
        fill=color,
        line=None,
        size=size,
        color=WHITE,
        bold=True,
        margins=(0.0, 0.0, 0.0, 0.0),
    )


def add_dot(slide, x: float, y: float, color: str, diameter: float = 0.10):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter)
    )
    set_fill(shape, color)
    set_line(shape, None)
    return shape


def add_chevron(slide, x: float, y: float, w: float = 0.24, h: float = 0.30, color: str = CYAN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, color)
    set_line(shape, None)
    return shape


def add_connector(slide, x1: float, y1: float, x2: float, y2: float, color: str = LINE, width: float = 1.0):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    set_line(shape, color, width)
    return shape


def clone_chrome_slide(prs: Presentation, source_slide):
    """Clone the source slide, keeping only its common header/footer chrome."""
    slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in source_slide.shapes:
        slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    # The original deck's common chrome is the first fifteen shapes on content slides.
    for shape in list(slide.shapes)[15:]:
        element = shape._element
        element.getparent().remove(element)
    return slide


def set_header(slide, section: str, title: str, subtitle: str) -> None:
    slide.shapes[1].text_frame.paragraphs[0].runs[0].text = section
    slide.shapes[2].text_frame.paragraphs[0].runs[0].text = title
    slide.shapes[3].text_frame.paragraphs[0].runs[0].text = subtitle


def add_slide_7(slide) -> None:
    set_header(
        slide,
        "YOLO FINE-TUNING LOOP",
        "从基线到发布：YOLO 微调优化闭环",
        "固定数据划分与评估口径，逐项验证数据、结构与推理优化，确保收益可复现、可归因。",
    )

    add_box(slide, 0.58, 1.48, 12.14, 2.20, fill=NAVY, line=NAVY)
    add_text(slide, 0.84, 1.63, 2.00, 0.22, "单变量推进｜5 个关口", size=9.0, color=CYAN, bold=True)

    items = [
        ("01", "数据审计", "类别 / 框尺度\n重复与划分泄漏", BLUE),
        ("02", "基线复现", "预训练权重\n固定 split / seed", CYAN),
        ("03", "单项改动", "分辨率 → P2\n→ CA → 切片", ORANGE),
        ("04", "消融验证", "AP_small / Recall\nFP / latency", PURPLE),
        ("05", "发布门槛", "精度 + 速度\n体积 + 可追溯", GREEN),
    ]
    x_positions = [0.86, 3.18, 5.50, 7.82, 10.14]
    for index, ((number, heading, body, accent), x) in enumerate(zip(items, x_positions)):
        add_box(slide, x, 1.96, 1.91, 1.33, fill=NAVY_2, line="183D5B")
        add_badge(slide, x + 0.14, 2.10, number, accent, 8.0)
        add_text(slide, x + 0.58, 2.08, 1.14, 0.26, heading, size=10.5, color=WHITE, bold=True)
        add_text(slide, x + 0.14, 2.52, 1.60, 0.50, body, size=8.4, color=LIGHT_MUTED, bold=False, align=PP_ALIGN.CENTER)
        if index < len(items) - 1:
            add_chevron(slide, x + 1.99, 2.48, 0.22, 0.28, CYAN)

    add_box(slide, 0.58, 3.92, 7.42, 2.68, fill=WHITE, line=LINE)
    add_text(slide, 0.84, 4.14, 2.60, 0.25, "三段式训练节奏", size=11.5, color=BLUE, bold=True)
    phases = [
        ("WARM-UP", "0–3 epoch", "稳定学习率与梯度\n小数据可选冻结骨干", BLUE, PALE_BLUE),
        ("MAIN", "主体训练", "全网络 + cos_lr\npatience=15 早停", CYAN, PALE_CYAN),
        ("CALIBRATE", "最后 10 epoch", "close_mosaic=10\n收敛到真实图像分布", GREEN, PALE_GREEN),
    ]
    phase_x = [0.84, 3.18, 5.52]
    for idx, ((tag, epoch, body, accent, pale), x) in enumerate(zip(phases, phase_x)):
        add_box(slide, x, 4.58, 1.98, 1.55, fill=pale, line=LINE)
        add_text(slide, x + 0.15, 4.72, 1.68, 0.18, tag, size=8.0, color=accent, bold=True)
        add_text(slide, x + 0.15, 4.96, 1.68, 0.25, epoch, size=11.2, color=TEXT, bold=True)
        add_text(slide, x + 0.15, 5.35, 1.68, 0.52, body, size=8.3, color=MUTED, bold=False)
        if idx < 2:
            add_chevron(slide, x + 2.04, 5.17, 0.22, 0.28, BLUE)

    add_box(slide, 8.24, 3.92, 4.48, 2.68, fill=PALE_GREEN, line=LINE)
    add_text(slide, 8.50, 4.14, 2.80, 0.25, "必须固化的实验指纹", size=11.5, color=GREEN, bold=True)
    fingerprint = [
        ("数据", "data.yaml、划分清单与数据哈希"),
        ("模型", "初始 .pt / 模型 YAML / 类别顺序"),
        ("参数", "args.yaml、随机种子、运行环境"),
        ("证据", "逐类指标、曲线、混淆矩阵与权重"),
    ]
    y = 4.63
    for label, body in fingerprint:
        add_box(slide, 8.50, y, 0.58, 0.30, label, fill=GREEN, line=None, size=8.2, color=WHITE, bold=True)
        add_text(slide, 9.22, y - 0.01, 3.05, 0.32, body, size=8.6, color=TEXT, bold=(label == "证据"))
        y += 0.43

    add_box(
        slide,
        0.58,
        6.74,
        12.14,
        0.30,
        "决策顺序：先复现 baseline，再一次只改一个变量；验证集用于选择，独立测试集只用于最终报告。",
        fill=PALE_ORANGE,
        line=None,
        size=8.8,
        color=ORANGE,
        bold=True,
    )


def add_tile_grid(slide, x: float, y: float, w: float, h: float) -> None:
    add_box(slide, x, y, w, h, fill="113553", line="2C5978", radius=False)
    # Tower-like silhouette and a tiny target near a tile boundary.
    add_connector(slide, x + 0.44, y + 1.40, x + 0.88, y + 0.25, LIGHT_MUTED, 1.4)
    add_connector(slide, x + 1.35, y + 1.40, x + 0.88, y + 0.25, LIGHT_MUTED, 1.4)
    add_connector(slide, x + 0.58, y + 0.96, x + 1.22, y + 0.96, LIGHT_MUTED, 1.0)
    add_connector(slide, x + 0.69, y + 0.66, x + 1.10, y + 0.66, LIGHT_MUTED, 1.0)
    for fraction in (0.38, 0.70):
        add_connector(slide, x + w * fraction, y, x + w * fraction, y + h, CYAN, 0.8)
    add_connector(slide, x, y + h * 0.53, x + w, y + h * 0.53, CYAN, 0.8)
    add_box(slide, x + 1.46, y + 0.70, 0.38, 0.24, fill=PALE_RED, line=RED, radius=False)
    add_text(slide, x + 1.50, y + 0.72, 0.30, 0.16, "tiny", size=5.8, color=RED, bold=True, align=PP_ALIGN.CENTER)


def add_slide_8(slide) -> None:
    set_header(
        slide,
        "SMALL-OBJECT DATA PIPELINE",
        "小目标优化（一）：尺度分析与重叠切片",
        "先用标注尺度决定是否引入高分辨率分支，再用切片放大小目标；恢复原图坐标后消除重叠检测。",
    )

    add_box(slide, 0.58, 1.48, 3.58, 5.15, fill=WHITE, line=LINE)
    add_text(slide, 0.84, 1.72, 2.86, 0.24, "01｜先判断“是否真小”", size=11.5, color=BLUE, bold=True)
    add_text(slide, 0.84, 2.10, 2.90, 0.42, "analyse_yolo_dataset 扫描标注框，输出尺度分布与结构建议。", size=8.9, color=MUTED)

    stat_cards = [
        ("判定阈值", "input_size / 64", "640 输入 → 10 px", BLUE, PALE_BLUE),
        ("核心统计", "短边中位数", "阈值以下的占比", CYAN, PALE_CYAN),
        ("形状聚类", "IoU K-means", "6 组宽高原型", PURPLE, PALE_PURPLE),
    ]
    y = 2.72
    for title, value, note, accent, pale in stat_cards:
        add_box(slide, 0.84, y, 2.98, 0.78, fill=pale, line=LINE)
        add_text(slide, 1.02, y + 0.12, 0.78, 0.20, title, size=8.2, color=accent, bold=True)
        add_text(slide, 1.86, y + 0.08, 1.70, 0.25, value, size=10.6, color=TEXT, bold=True)
        add_text(slide, 1.02, y + 0.43, 2.54, 0.18, note, size=7.8, color=MUTED)
        y += 0.93

    add_box(slide, 0.84, 5.57, 2.98, 0.70, fill=PALE_ORANGE, line=None)
    add_text(slide, 1.02, 5.66, 2.60, 0.44, "P2 建议：短边中位数 < 阈值\n或阈值以下占比 ≥ 20%", size=8.7, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.84, 6.31, 2.98, 0.17, "尺度需按实际 letterbox / 推理分辨率复核。", size=7.3, color=RED, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, 4.40, 1.48, 8.32, 3.42, fill=NAVY, line=NAVY)
    add_text(slide, 4.68, 1.70, 3.20, 0.24, "02｜重叠切片如何落地", size=11.5, color=CYAN, bold=True)
    add_tile_grid(slide, 4.72, 2.10, 2.28, 1.62)
    add_text(slide, 4.72, 3.84, 2.28, 0.30, "高分辨率原图｜边界目标被相邻切片共同观察", size=7.4, color=LIGHT_MUTED, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "生成窗口", "tile 1024\noverlap 20%", BLUE),
        ("2", "局部推理", "目标像素\n占比提高", CYAN),
        ("3", "坐标恢复", "bbox +\n(x, y) 偏移", ORANGE),
        ("4", "同类融合", "IoU 0.55\n置信度加权", GREEN),
    ]
    xs = [7.28, 8.63, 9.98, 11.33]
    for idx, ((number, heading, note, accent), x) in enumerate(zip(steps, xs)):
        add_box(slide, x, 2.12, 1.10, 1.70, fill=NAVY_2, line="183D5B")
        add_badge(slide, x + 0.37, 2.28, number, accent, 8.0)
        add_text(slide, x + 0.08, 2.76, 0.94, 0.24, heading, size=8.7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.08, 3.10, 0.94, 0.48, note, size=7.4, color=LIGHT_MUTED, align=PP_ALIGN.CENTER)
        if idx < 3:
            add_chevron(slide, x + 1.13, 2.84, 0.17, 0.24, CYAN)

    add_box(slide, 4.40, 5.12, 4.02, 1.51, fill=PALE_BLUE, line=LINE)
    add_text(slide, 4.68, 5.34, 2.50, 0.22, "WEB：全图 + 切片", size=10.8, color=BLUE, bold=True)
    add_dot(slide, 4.70, 5.80, BLUE)
    add_text(slide, 4.92, 5.72, 3.08, 0.32, "高分辨率图像先全图检测，再叠加裁剪结果", size=8.2, color=TEXT)
    add_dot(slide, 4.70, 6.18, BLUE)
    add_text(slide, 4.92, 6.10, 3.08, 0.32, "类别感知 NMS 合并，避免跨类别误抑制", size=8.2, color=TEXT, bold=True)

    add_box(slide, 8.66, 5.12, 4.06, 1.51, fill=PALE_GREEN, line=LINE)
    add_text(slide, 8.94, 5.34, 2.60, 0.22, "ROS：按需启用切片", size=10.8, color=GREEN, bold=True)
    add_dot(slide, 8.96, 5.80, GREEN)
    add_text(slide, 9.18, 5.72, 3.10, 0.32, "detector.yaml 默认关闭，远距场景再开启", size=8.2, color=TEXT)
    add_dot(slide, 8.96, 6.18, GREEN)
    add_text(slide, 9.18, 6.10, 3.10, 0.32, "同类 WBF：加权坐标，保留组内最高置信度", size=8.2, color=TEXT, bold=True)

    add_box(
        slide,
        0.58,
        6.75,
        12.14,
        0.28,
        "收益来自“增加目标有效像素”；代价是切片数带来的延迟，因此要同时记录召回提升、FP / 图与 p95 延迟。",
        fill=PALE_RED,
        line=None,
        size=8.4,
        color=RED,
        bold=True,
    )


def add_slide_9(slide) -> None:
    set_header(
        slide,
        "SMALL-OBJECT MODEL DESIGN",
        "小目标优化（二）：P2 + Coordinate Attention",
        "将 stride=4 的浅层高分辨率特征送入检测头，并用方向敏感注意力增强细长目标与缺片位置响应。",
    )

    add_box(slide, 0.58, 1.48, 4.12, 5.17, fill=PALE_BLUE, line=LINE)
    add_text(slide, 0.84, 1.72, 2.85, 0.24, "P2：保留更细的空间网格", size=11.5, color=BLUE, bold=True)
    add_text(slide, 0.84, 2.06, 3.42, 0.32, "640 输入下，P2 特征图为 160×160；比 P3 多 4 倍网格位置。", size=8.4, color=MUTED)
    levels = [
        ("P2 / 4", "160 × 160", 2.36, 2.22, CYAN, NAVY),
        ("P3 / 8", "80 × 80", 3.16, 1.95, BLUE, WHITE),
        ("P4 / 16", "40 × 40", 3.96, 1.68, PURPLE, WHITE),
        ("P5 / 32", "20 × 20", 4.76, 1.41, ORANGE, WHITE),
    ]
    for label, size_text, y, width, accent, fill in levels:
        text_color = WHITE if fill == NAVY else TEXT
        add_box(slide, 0.90, y, width, 0.56, fill=fill, line=accent, radius=False)
        add_text(slide, 1.05, y + 0.07, 0.82, 0.20, label, size=9.6, color=accent, bold=True)
        add_text(slide, 1.92, y + 0.07, max(0.55, width - 1.10), 0.20, size_text, size=8.0, color=text_color, bold=True, align=PP_ALIGN.RIGHT)
        add_connector(slide, 0.90 + width, y + 0.28, 3.55, y + 0.28, accent, 1.1)
    add_box(slide, 3.52, 2.34, 0.78, 3.00, "v10\nDetect\n\n4-scale", fill=NAVY, line=NAVY, size=9.0, color=WHITE, bold=True)
    add_box(slide, 0.88, 5.70, 3.52, 0.60, "检测头由 P3/P4/P5 扩为 P2/P3/P4/P5", fill=WHITE, line=LINE, size=8.8, color=BLUE, bold=True)

    add_box(slide, 4.96, 1.48, 4.16, 5.17, fill=NAVY, line=NAVY)
    add_text(slide, 5.24, 1.72, 3.20, 0.24, "Coordinate Attention：保留方向坐标", size=11.2, color=CYAN, bold=True)
    add_text(slide, 5.24, 2.06, 3.54, 0.36, "对 H、W 两个方向分别池化，生成位置敏感的通道权重。", size=8.3, color=LIGHT_MUTED)
    add_box(slide, 5.38, 2.65, 1.05, 0.55, "X\nC×H×W", fill=NAVY_2, line="2C5978", size=8.3, color=WHITE, bold=True)
    add_chevron(slide, 6.53, 2.77, 0.20, 0.28, CYAN)
    add_box(slide, 6.82, 2.50, 1.78, 0.42, "沿 W 平均池化 → H", fill=PALE_CYAN, line=None, size=7.8, color=TEXT, bold=True)
    add_box(slide, 6.82, 3.02, 1.78, 0.42, "沿 H 平均池化 → W", fill=PALE_BLUE, line=None, size=7.8, color=TEXT, bold=True)
    add_chevron(slide, 6.20, 3.62, 0.22, 0.30, CYAN)
    add_box(slide, 6.55, 3.57, 2.05, 0.70, "Concat → 1×1 Conv\nBN + HSwish", fill=NAVY_2, line="2C5978", size=8.2, color=WHITE, bold=True)
    add_chevron(slide, 6.20, 4.47, 0.22, 0.30, CYAN)
    add_box(slide, 5.42, 4.48, 1.36, 0.65, "sigmoid\naₕ", fill=PALE_GREEN, line=None, size=8.5, color=GREEN, bold=True)
    add_box(slide, 7.22, 4.48, 1.36, 0.65, "sigmoid\naᵥ", fill=PALE_PURPLE, line=None, size=8.5, color=PURPLE, bold=True)
    add_connector(slide, 6.10, 5.13, 6.70, 5.48, GREEN, 1.0)
    add_connector(slide, 7.90, 5.13, 7.30, 5.48, PURPLE, 1.0)
    add_box(slide, 6.03, 5.45, 1.96, 0.52, "Y = X × aₕ × aᵥ", fill=CYAN, line=None, size=9.1, color=NAVY, bold=True)
    add_text(slide, 5.38, 6.12, 3.30, 0.23, "输出形状不变，可插入现有 YAML 解析链路", size=7.9, color=LIGHT_MUTED, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, 9.36, 1.48, 3.36, 5.17, fill=WHITE, line=LINE)
    add_text(slide, 9.62, 1.72, 2.52, 0.24, "代码落点与验证方式", size=11.2, color=GREEN, bold=True)
    cards = [
        ("MODEL YAML", "P2 / P3 后插入 CoordAtt\n检测头输出 [P2…P5]", BLUE, PALE_BLUE),
        ("PYTORCH", "hidden=max(8, C/32)\n惰性构建，保持张量形状", CYAN, PALE_CYAN),
        ("ABLATION", "B0 baseline → B1 P2\n→ B2 P2+CA", PURPLE, PALE_PURPLE),
        ("COST", "同时记录 VRAM、FLOPs\n吞吐与 p95 延迟", ORANGE, PALE_ORANGE),
    ]
    y = 2.18
    for tag, body, accent, pale in cards:
        add_box(slide, 9.62, y, 2.84, 0.88, fill=pale, line=LINE)
        add_text(slide, 9.80, y + 0.10, 0.88, 0.18, tag, size=7.4, color=accent, bold=True)
        add_text(slide, 9.80, y + 0.34, 2.42, 0.38, body, size=8.1, color=TEXT, bold=(tag == "ABLATION"))
        y += 1.00

    add_box(
        slide,
        0.58,
        6.76,
        12.14,
        0.27,
        "YOLOv10 仍为 anchor-free：IoU 宽高聚类只用于选择 P2 / TAL / 增广方案，不会生成 anchors。",
        fill=PALE_ORANGE,
        line=None,
        size=8.4,
        color=ORANGE,
        bold=True,
    )


def add_slide_10(slide) -> None:
    set_header(
        slide,
        "TRAINING RECIPE & ABLATION",
        "微调配方、消融实验与上线判据",
        "训练参数先作为可复现起点；最终方案由固定划分的精度证据和同硬件性能预算共同决定。",
    )

    add_box(slide, 0.58, 1.48, 5.12, 5.16, fill=NAVY, line=NAVY)
    add_text(slide, 0.86, 1.72, 2.75, 0.25, "当前可复现训练入口", size=11.5, color=CYAN, bold=True)
    add_box(slide, 0.86, 2.10, 4.56, 1.31, fill=NAVY_2, line="183D5B")
    command = (
        "python -m insulens_perception.train\n"
        "  --small-object-model --epochs 60\n"
        "  --imgsz 640 --batch 16\n"
        "  --tal-topk 10 --patience 15"
    )
    add_text(slide, 1.05, 2.23, 4.15, 1.00, command, size=8.2, color=WHITE, bold=False, font=MONO, valign=MSO_ANCHOR.TOP)

    add_text(slide, 0.86, 3.68, 2.20, 0.22, "训练脚本中的固定起点", size=9.5, color=LIGHT_MUTED, bold=True)
    settings = [
        ("优化", "optimizer=auto · cos_lr=True", BLUE),
        ("稳定", "patience=15 · AMP=False", CYAN),
        ("几何", "degrees=12 · translate=.15 · scale=.45", ORANGE),
        ("混合", "fliplr=.5 · mixup=.05", PURPLE),
        ("收尾", "close_mosaic=10", GREEN),
    ]
    y = 4.02
    for label, body, accent in settings:
        add_box(slide, 0.88, y, 0.68, 0.31, label, fill=accent, line=None, size=8.0, color=WHITE, bold=True)
        add_text(slide, 1.72, y - 0.01, 3.50, 0.32, body, size=8.5, color=WHITE, bold=(label == "收尾"))
        y += 0.43
    add_box(slide, 0.86, 6.22, 4.56, 0.25, "预训练权重仅在 --model 指向 .pt 时启用", fill=PALE_CYAN, line=None, size=7.8, color=NAVY, bold=True)

    add_box(slide, 5.96, 1.48, 6.76, 3.62, fill=WHITE, line=LINE)
    add_text(slide, 6.22, 1.70, 3.25, 0.24, "最小消融矩阵｜每组固定 split / seed", size=11.2, color=BLUE, bold=True)
    col_x = [6.22, 6.84, 9.12]
    col_w = [0.60, 2.26, 3.32]
    headers = ["组", "只改变什么", "必须比较"]
    for x, w, label in zip(col_x, col_w, headers):
        add_box(slide, x, 2.08, w, 0.40, label, fill=NAVY, line=NAVY, radius=False, size=8.2, color=WHITE, bold=True)
    rows = [
        ("B0", "YOLOv10s @ 640", "基线：逐类 AP / Recall / FP"),
        ("B1", "提高 imgsz 或切片", "像素收益 vs p95 latency"),
        ("B2", "+ P2 检测分支", "AP_small vs VRAM / FLOPs"),
        ("B3", "+ Coordinate Attention", "定位收益 vs 吞吐"),
        ("B4", "P2 + CA + 切片", "组合收益、误报与稳定性"),
    ]
    y = 2.48
    for idx, row in enumerate(rows):
        fill = WHITE if idx % 2 == 0 else PALE_GRAY
        for x, w, value in zip(col_x, col_w, row):
            add_box(slide, x, y, w, 0.45, value, fill=fill, line=LINE, radius=False, size=7.8, color=(BLUE if x == col_x[0] else TEXT), bold=(x != col_x[2]))
        y += 0.45
    add_text(slide, 6.22, 4.80, 6.02, 0.17, "资源允许时重复 ≥3 次，报告均值与波动；不要只展示最好的一次。", size=7.6, color=RED, bold=True, align=PP_ALIGN.RIGHT)

    add_box(slide, 5.96, 5.32, 6.76, 1.32, fill=PALE_GREEN, line=LINE)
    add_text(slide, 6.22, 5.51, 1.70, 0.21, "上线判据｜同时满足", size=10.4, color=GREEN, bold=True)
    gates = [
        ("AP_small", BLUE),
        ("逐类 Recall", CYAN),
        ("FP / 图", RED),
        ("p95 延迟", ORANGE),
        ("VRAM / 体积", PURPLE),
    ]
    x = 6.22
    for label, accent in gates:
        add_box(slide, x, 5.88, 1.05, 0.36, label, fill=WHITE, line=accent, size=7.9, color=accent, bold=True)
        x += 1.17
    add_text(slide, 6.22, 6.34, 5.92, 0.18, "优先保证缺陷召回；精度收益必须落在部署时延与显存预算内。", size=7.8, color=TEXT, bold=True, align=PP_ALIGN.CENTER)

    add_box(
        slide,
        0.58,
        6.76,
        12.14,
        0.27,
        "实现边界：--tal-topk 当前仅写入 small_object_experiment.json；分配器仍由已安装的 Ultralytics 版本管理。",
        fill=PALE_ORANGE,
        line=None,
        size=8.3,
        color=ORANGE,
        bold=True,
    )


def update_page_numbers(prs: Presentation) -> None:
    for number, slide in enumerate(prs.slides, start=1):
        if number == 1:
            continue
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            x = shape.left / 914400
            y = shape.top / 914400
            text = shape.text.strip()
            if x > 12.0 and y < 0.75 and re.fullmatch(r"\d{2}", text):
                shape.text_frame.paragraphs[0].runs[0].text = f"{number:02d}"
            elif x > 11.5 and y > 7.0 and re.fullmatch(r"\d+\s*/\s*\d+", text):
                shape.text_frame.paragraphs[0].runs[0].text = f"{number} / 12"


def main() -> None:
    prs = Presentation(SOURCE)
    original_slide_7 = prs.slides[6]
    chrome_source = prs.slides[7]

    added = [clone_chrome_slide(prs, chrome_source) for _ in range(4)]
    add_slide_7(added[0])
    add_slide_8(added[1])
    add_slide_9(added[2])
    add_slide_10(added[3])

    # Reorder: original slides 1–6, new slides 7–10, original slides 7–8.
    slide_id_list = prs.slides._sldIdLst
    insertion_index = list(prs.slides).index(original_slide_7)
    for slide in added:
        slide_id = None
        # python-pptx exposes the slide id through its relationship. Locate it by part.
        for candidate in list(slide_id_list):
            related = prs.part.related_part(candidate.rId)
            if related is slide.part:
                slide_id = candidate
                break
        if slide_id is None:
            raise RuntimeError("Unable to locate newly added slide id")
        slide_id_list.remove(slide_id)
        slide_id_list.insert(insertion_index, slide_id)
        insertion_index += 1

    update_page_numbers(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
